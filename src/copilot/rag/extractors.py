"""Document extractors for mixed-format corpora.

Each extractor returns a list of `Block`s — a unit of text with provenance
(page number and/or section heading). Downstream chunking groups blocks into
appropriately-sized chunks while preserving that provenance in metadata, so a
retrieved chunk can cite "source X, p. 42".

Supported:
  .pdf         → text layer per page; if a page has no text layer, OCR it
  .docx        → paragraphs + headings (headings become section markers)
  .pptx        → one block per slide (slide number as "page")
  .txt / .md   → whole file, markdown headings become section markers
  .jpg/.jpeg/.png → OCR

OCR uses pytesseract (needs the `tesseract-ocr` system binary, plus the
`deu` + `eng` language packs for this bilingual corpus).
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field
from pathlib import Path

# --- optional heavy deps imported lazily so import never hard-fails ----------

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff"}


@dataclass
class Block:
    text: str
    page: int | None = None          # 1-based page/slide number, if applicable
    section: str | None = None       # nearest heading, if known
    meta: dict = field(default_factory=dict)


# --------------------------------------------------------------------------- #
# PDF (with OCR fallback per page)
# --------------------------------------------------------------------------- #
def extract_pdf(path: Path, ocr_lang: str = "deu+eng", ocr_threshold: int = 20) -> list[Block]:
    """Extract text per page. Pages whose text layer is shorter than
    `ocr_threshold` chars are treated as scanned and OCR'd."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    blocks: list[Block] = []
    scanned_pages: list[int] = []

    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if len(text) < ocr_threshold:
            scanned_pages.append(i)
            text = _ocr_pdf_page(path, i - 1, ocr_lang) or text
        if text:
            blocks.append(Block(text=text, page=i, meta={"ocr": i in scanned_pages}))

    return blocks


def _ocr_pdf_page(path: Path, page_index: int, lang: str) -> str:
    """Render a single PDF page to an image and OCR it. Best-effort; returns
    empty string if the OCR stack isn't available."""
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except Exception:
        return ""
    try:
        images = convert_from_path(
            str(path), first_page=page_index + 1, last_page=page_index + 1, dpi=200
        )
        if not images:
            return ""
        return pytesseract.image_to_string(images[0], lang=lang).strip()
    except Exception:
        return ""


# --------------------------------------------------------------------------- #
# DOCX
# --------------------------------------------------------------------------- #
def extract_docx(path: Path) -> list[Block]:
    from docx import Document

    doc = Document(str(path))
    blocks: list[Block] = []
    current_section: str | None = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading") or style == "title":
            current_section = text
            blocks.append(Block(text=text, section=current_section, meta={"heading": True}))
        else:
            blocks.append(Block(text=text, section=current_section))

    # Tables → rows joined; keeps rubric-style content intact.
    for table in doc.tables:
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells)
            for row in table.rows
        ]
        table_text = "\n".join(r for r in rows if r.strip(" |"))
        if table_text:
            blocks.append(Block(text=table_text, section=current_section, meta={"table": True}))

    return blocks


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #
def extract_pptx(path: Path) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[Block] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            blocks.append(Block(text="\n".join(parts), page=i, meta={"slide": True}))
    return blocks


# --------------------------------------------------------------------------- #
# Plain text / markdown
# --------------------------------------------------------------------------- #
def extract_text(path: Path) -> list[Block]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    blocks: list[Block] = []
    current_section: str | None = None
    buf: list[str] = []

    def flush():
        if buf:
            blocks.append(Block(text="\n".join(buf).strip(), section=current_section))
            buf.clear()

    for line in raw.splitlines():
        if line.lstrip().startswith("#"):        # markdown heading
            flush()
            current_section = line.lstrip("#").strip()
            blocks.append(Block(text=current_section, section=current_section, meta={"heading": True}))
        else:
            buf.append(line)
    flush()
    return [b for b in blocks if b.text]


# --------------------------------------------------------------------------- #
# Image OCR
# --------------------------------------------------------------------------- #
def extract_image(path: Path, ocr_lang: str = "deu+eng") -> list[Block]:
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        return []
    try:
        text = pytesseract.image_to_string(Image.open(path), lang=ocr_lang).strip()
    except Exception:
        return []
    return [Block(text=text, meta={"ocr": True})] if text else []


# --------------------------------------------------------------------------- #
# Dispatch
# --------------------------------------------------------------------------- #
def extract(path: Path, ocr_lang: str = "deu+eng") -> list[Block]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path, ocr_lang=ocr_lang)
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext in {".txt", ".md"}:
        return extract_text(path)
    if ext in IMAGE_EXTS:
        return extract_image(path, ocr_lang=ocr_lang)
    return []


SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"} | IMAGE_EXTS
