#!/usr/bin/env python3
"""
Bilingual (EN/DE) document -> chunk pipeline.
Run: source .venv/bin/activate && python -u pipeline.py
Override: INPUT_DIR=./sample OUTPUT_DIR=./output_sample python -u pipeline.py
"""

import os, sys, json, shutil, hashlib, mimetypes, unicodedata
from pathlib import Path
from datetime import datetime, timezone
from collections import Counter

import regex as re                       # variable-width lookbehind; stdlib `re` can't
import numpy as np
import fitz
fitz.TOOLS.mupdf_display_errors(False)   # benign ICC/stream warnings

import ftfy
from docx import Document as Docx
from pptx import Presentation
from PIL import Image
from PIL.ExifTags import TAGS
import pytesseract
from langdetect import detect_langs, DetectorFactory
import tiktoken
from sentence_transformers import SentenceTransformer

DetectorFactory.seed = 0

# ── config ────────────────────────────────────────────────────────────
INPUT_DIR  = os.environ.get("INPUT_DIR",  "./raw/documents")
OUTPUT_DIR = os.environ.get("OUTPUT_DIR", "./output")
TARGET_TOKENS, MAX_TOKENS, MIN_TOKENS, OVERLAP_TOKENS = 400, 700, 80, 60
BREAKPOINT_PERCENTILE = 85
MIN_DOC_TOKENS = 50
OCR_LANGS, OCR_DPI = "deu+eng", 300

OCR_AVAILABLE = shutil.which("tesseract") is not None
if not OCR_AVAILABLE:
    print("WARNING: tesseract not found — scanned pages yield no text", file=sys.stderr)

ENC = tiktoken.get_encoding("cl100k_base")
ntok = lambda s: len(ENC.encode(s, disallowed_special=()))

print("loading embedder…", file=sys.stderr)
EMBEDDER = SentenceTransformer("intfloat/multilingual-e5-base")


# ── preprocessing ─────────────────────────────────────────────────────
GERMAN_FIXES = {
    "ﬀ": "ff", "ﬁ": "fi", "ﬂ": "fl", "ﬃ": "ffi", "ﬄ": "ffl", "ﬅ": "ft", "ﬆ": "st",
    "„": '"', "“": '"', "”": '"', "‚": "'", "‘": "'", "’": "'",
    "–": "-", "—": "-", "−": "-", "\u00ad": "",
    "\u200b": "", "\u200c": "", "\u200d": "", "\ufeff": "", "\u00a0": " ",
    "\u202f": " ", "\u2009": " ", "\u2007": " ",
}
OCR_UMLAUT = [(r"A\u0308", "Ä"), (r"O\u0308", "Ö"), (r"U\u0308", "Ü"),
              (r"a\u0308", "ä"), (r"o\u0308", "ö"), (r"u\u0308", "ü")]

# tesseract reads capital I as lowercase l — fatal for a corpus about "KI"
OCR_CONFUSIONS = [
    (r"\bKl\b", "KI"),
    (r"\bKl(?=[-\u2010])", "KI"),
    (r"(?<=Klasse |Kategorie |Stufe |Level )Il\b", "II"),
    (r"(?<=Klasse |Kategorie |Stufe |Level )lll\b", "III"),
    (r"\bAl Act\b", "AI Act"),
    (r"\bAl-\b", "AI-"),
]

DE_ABBREV = r"(?:z\.\s?B|d\.\s?h|u\.\s?a|bzw|ggf|inkl|exkl|Nr|Abs|Art|Ziff|Bsp|evtl|ca|vgl|s\.\s?o|s\.\s?u|Dr|Prof|Dipl|Ing|Mio|Mrd|Str|bspw|i\.\s?d\.\s?R|u\.\s?U|o\.\s?g|sog|St|Hr|Fr)"
EN_ABBREV = r"(?:e\.g|i\.e|etc|vs|cf|approx|Fig|Eq|No|Dr|Prof|Mr|Mrs|Ms|Inc|Ltd|Corp|St|Jr|Sr|al)"

SENT_SPLIT = re.compile(
    rf"(?<!\b{DE_ABBREV})(?<!\b{EN_ABBREV})"
    r"(?<!\b[A-ZÄÖÜ])(?<!\b\d)"
    r"(?<=[.!?:;])[\"'»«”’)\]]*\s+(?=[A-ZÄÖÜ„«\"'\d])"
)

NOISE_PATTERNS = [
    re.compile(r"^\s*(?:seite|page|s\.)\s*[-–]?\s*\d+\s*(?:(?:von|of|/)\s*\d+)?\s*$", re.I | re.M),
    re.compile(r"^\s*[-–—]\s*\d+\s*[-–—]\s*$", re.M),
    re.compile(r"^\s*\d{1,3}\s*$", re.M),
    re.compile(r"^\s*(?:vertraulich|confidential|internal|entwurf|draft)\s*$", re.I | re.M),
    re.compile(r"(?:Stand|Version|Rev\.?|Last updated)[:\s]+\d{1,2}[./]\d{1,2}[./]\d{2,4}\s*$", re.I | re.M),
    re.compile(r"^\s*(?:https?://\S+|www\.\S+)\s*$", re.M),
    re.compile(r"\.{4,}\s*\d+\s*$", re.M),
]

# narrow: must not catch "RPO_Bachelor_Master_WiSe" (a real regulation)
EXCLUDE = re.compile(r"(?:ppt-?master|word-?template|vorlage|_muster\b|^muster)", re.I)


def norm_key(t):
    return re.sub(r"\W+", "", t.lower())[:200]


def fix_encoding(t):
    t = ftfy.fix_text(t, normalization="NFC")
    for a, b in GERMAN_FIXES.items():
        t = t.replace(a, b)
    for p, r in OCR_UMLAUT:
        t = re.sub(p, r, t)
    return unicodedata.normalize("NFC", t)


def fix_ocr_confusions(t):
    for p, r in OCR_CONFUSIONS:
        t = re.sub(p, r, t)
    return t


def repair_ocr_spacing(t):
    return re.sub(r"\b(?:[A-ZÄÖÜ]\s){3,}[A-ZÄÖÜ]\b", lambda m: m.group(0).replace(" ", ""), t)


def dehyphenate(t):
    t = re.sub(r"(?<=[a-zäöüß])[-\u2010\u00ad]\s*\n\s*(?=[a-zäöüß])", "", t)   # Betriebs-\nrat
    t = re.sub(r"(?<=[a-zäöüß])[-\u2010]\s*\n\s*(?=[A-ZÄÖÜ])", "-", t)         # Arbeits-\nRecht
    t = re.sub(r"(?<=[a-zäöüß,;])\n(?=[a-zäöüß])", " ", t)                     # soft wrap
    return t


def strip_noise(t):
    for p in NOISE_PATTERNS:
        t = p.sub("", t)
    t = re.sub(r"[ \t\u00a0]{2,}", " ", t)
    t = re.sub(r" +\n", "\n", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


def preprocess(raw, ocr=False):
    t = fix_encoding(raw)
    if ocr:
        t = fix_ocr_confusions(t)    # only OCR text — clean docs may contain a real "Kl"
    return strip_noise(dehyphenate(repair_ocr_spacing(t)))


def drop_boilerplate(pages, thresh=0.55):
    if len(pages) < 3:
        return pages
    top, bot = Counter(), Counter()
    for p in pages:
        lines = [l.strip() for l in p["text"].split("\n") if l.strip()]
        if lines:
            top.update(lines[:2]); bot.update(lines[-2:])
    n = len(pages)
    bad = {l for c in (top, bot) for l, k in c.items()
           if k / n > thresh and 3 < len(l) < 130 and not re.search(r"[.!?]$", l)}
    for p in pages:
        p["text"] = "\n".join(l for l in p["text"].split("\n") if l.strip() not in bad)
    return pages


def detect_lang(t):
    s = t[:3000]
    if len(s.strip()) < 25:
        return "unknown", 0.0
    try:
        r = detect_langs(s)[0]
        return r.lang, round(r.prob, 3)
    except Exception:
        return "unknown", 0.0


def doc_is_empty(t):
    """Document-level check. Deliberately NOT is_garbage(): that function's
    ratio heuristics (single-char words, unique-word fraction) are tuned for
    ~400-token chunks and misfire badly on book-length text, where vocabulary
    naturally saturates."""
    if ntok(t) < MIN_DOC_TOKENS:
        return True
    if len(t) and sum(c.isalpha() for c in t) / len(t) < 0.40:
        return True
    return False


def is_garbage(t):
    if len(t) < 30:
        return True
    if sum(c.isalpha() for c in t) / len(t) < 0.55:
        return True
    w = t.split()
    if not w:
        return True
    if sum(len(x) == 1 for x in w) / len(w) > 0.4:
        return True
    if len(set(w)) / len(w) < 0.25:
        return True
    return False


# ── extraction ────────────────────────────────────────────────────────
def _needs_ocr(page):
    return len(page.get_text("text").strip()) < 40


def _ocr_page(page):
    if not OCR_AVAILABLE:
        return ""
    pix = page.get_pixmap(dpi=OCR_DPI)
    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    return pytesseract.image_to_string(img, lang=OCR_LANGS, config="--psm 3")


def extract_pdf(p):
    doc = fitz.open(p)
    m = doc.metadata or {}
    pages, ocr_pages = [], 0
    for i, page in enumerate(doc):
        blocks = page.get_text("blocks", sort=True)      # reading order on 2-col layouts
        txt = "\n\n".join(b[4] for b in blocks if b[6] == 0 and b[4].strip())
        if _needs_ocr(page):
            txt = _ocr_page(page)
            ocr_pages += 1
        headings = [
            s["text"].strip()
            for b in page.get_text("dict")["blocks"] if b.get("type") == 0
            for l in b["lines"] for s in l["spans"]
            if s["size"] > 12.5 and (s["flags"] & 16) and len(s["text"].strip()) > 3
        ]
        pages.append({"page": i + 1, "text": txt, "headings": headings[:3]})
    toc = [t[1] for t in doc.get_toc()][:60]
    doc.close()
    return {"title": m.get("title"), "author": m.get("author"), "subject": m.get("subject"),
            "keywords": m.get("keywords"), "producer": m.get("producer"),
            "created": m.get("creationDate"), "modified": m.get("modDate"),
            "page_count": len(pages), "ocr_pages": ocr_pages,
            "is_scanned": ocr_pages > len(pages) * 0.5, "toc": toc}, pages


def extract_docx(p):
    d = Docx(p); cp = d.core_properties
    secs, cur, head = [], [], None
    for para in d.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        if para.style.name.startswith(("Heading", "Überschrift", "Title")):
            if cur:
                secs.append({"page": len(secs) + 1, "text": "\n".join(cur), "headings": [head] if head else []})
            cur, head = [t], t
        else:
            cur.append(t)
    for tbl in d.tables:
        rows = [" | ".join(c.text.strip() for c in r.cells) for r in tbl.rows]
        if rows:
            cur.append("[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]")
    if cur:
        secs.append({"page": len(secs) + 1, "text": "\n".join(cur), "headings": [head] if head else []})
    return {"title": cp.title, "author": cp.author, "subject": cp.subject, "keywords": cp.keywords,
            "created": str(cp.created), "modified": str(cp.modified),
            "page_count": len(secs), "ocr_pages": 0, "is_scanned": False}, secs


def extract_pptx(p):
    prs = Presentation(p); cp = prs.core_properties
    slides = []
    for i, s in enumerate(prs.slides):
        title, body = None, []
        for shp in s.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                if shp == s.shapes.title:
                    title = shp.text_frame.text.strip()
                else:
                    body.append(shp.text_frame.text)
            if getattr(shp, "has_table", False):
                body.append("[TABLE]\n" + "\n".join(
                    " | ".join(c.text for c in r.cells) for r in shp.table.rows) + "\n[/TABLE]")
        notes = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        txt = "\n".join(filter(None, [title, *body]))
        if notes:
            txt += f"\n[NOTES] {notes}"
        slides.append({"page": i + 1, "text": txt, "headings": [title] if title else []})
    return {"title": cp.title, "author": cp.author, "subject": cp.subject, "keywords": cp.keywords,
            "created": str(cp.created), "modified": str(cp.modified),
            "page_count": len(slides), "ocr_pages": 0, "is_scanned": False}, slides


def extract_image(p):
    img = Image.open(p)
    exif = {}
    try:
        exif = {TAGS.get(k, k): str(v) for k, v in (img._getexif() or {}).items()}
    except Exception:
        pass
    txt = ""
    if OCR_AVAILABLE:
        g = img.convert("L")
        if min(g.size) < 1000:
            g = g.resize((g.width * 2, g.height * 2), Image.LANCZOS)
        g = g.point(lambda x: 0 if x < 140 else 255)
        txt = pytesseract.image_to_string(g, lang=OCR_LANGS, config="--psm 3")
    return {"title": None, "author": exif.get("Artist"), "width": img.width, "height": img.height,
            "created": exif.get("DateTimeOriginal"), "camera": exif.get("Model"),
            "page_count": 1, "ocr_pages": 1, "is_scanned": True}, [{"page": 1, "text": txt, "headings": []}]


EXTRACTORS = {".pdf": extract_pdf, ".docx": extract_docx, ".pptx": extract_pptx,
              ".png": extract_image, ".jpg": extract_image, ".jpeg": extract_image,
              ".tif": extract_image, ".tiff": extract_image, ".webp": extract_image}


# ── chunking ──────────────────────────────────────────────────────────
def split_sentences(text):
    out = []
    for para in re.split(r"\n{2,}", text):
        para = para.strip()
        if not para:
            continue
        if para.startswith("[TABLE]"):
            out.append(para); continue
        for s in (x.strip() for x in SENT_SPLIT.split(para) if x.strip()):
            if ntok(s) <= MAX_TOKENS:
                out.append(s)
            else:                                    # 90-word German clause monsters
                out.extend(x.strip() for x in re.split(r"(?<=[,;:])\s+(?=\w)", s) if x.strip())
    return out


def _pack(sents, sims=None, thresh=None):
    chunks, cur, cur_tok = [], [], 0
    for i, s in enumerate(sents):
        st = ntok(s)
        brk = sims is not None and i > 0 and sims[i - 1] < thresh
        if cur and (cur_tok + st > MAX_TOKENS or (brk and cur_tok >= TARGET_TOKENS * 0.6)):
            chunks.append(" ".join(cur))
            back, bt = [], 0
            for prev in reversed(cur):               # sentence-level overlap
                pt = ntok(prev)
                if bt + pt > OVERLAP_TOKENS:
                    break
                back.insert(0, prev); bt += pt
            cur, cur_tok = back, bt
        cur.append(s); cur_tok += st
    if cur:
        chunks.append(" ".join(cur))
    merged = []
    for c in chunks:                                 # absorb runts
        if merged and ntok(c) < MIN_TOKENS and ntok(merged[-1]) + ntok(c) <= MAX_TOKENS:
            merged[-1] += " " + c
        else:
            merged.append(c)
    return merged


def semantic_chunk(text):
    sents = split_sentences(text)
    if not sents:
        return []
    if len(sents) < 4:
        return [" ".join(sents)]
    emb = EMBEDDER.encode([f"passage: {s}" for s in sents], normalize_embeddings=True,
                          show_progress_bar=False, batch_size=64)
    sims = np.einsum("ij,ij->i", emb[:-1], emb[1:])
    return _pack(sents, sims, np.percentile(sims, 100 - BREAKPOINT_PERCENTILE))


def chunk_document(dm, pages):
    chunks, idx = [], 0
    for pg in pages:
        text = preprocess(pg["text"], ocr=dm.get("is_scanned", False))
        if not text or is_garbage(text):
            continue
        for part in semantic_chunk(text):
            if is_garbage(part):
                continue
            lang, conf = detect_lang(part)
            head = pg["headings"][0] if pg["headings"] else None
            ctx = dm.get("title") or dm["file_name"]
            if head:
                ctx += f" > {head}"
            chunks.append({
                "chunk_id": f"{dm['doc_id'][:12]}_{idx:04d}", "doc_id": dm["doc_id"],
                "source_file": dm["file_name"], "file_type": dm["file_type"],
                "doc_title": dm.get("title") or dm["file_name"],
                "section_heading": head, "context_header": ctx,
                "embed_text": f"{ctx}\n\n{part}",     # embed this
                "text": part,                         # show this to the LLM
                "page": pg["page"], "chunk_index": idx,
                "language": lang, "language_conf": conf,
                "is_ocr": dm.get("is_scanned", False),
                "contains_table": "[TABLE]" in part,
                "char_count": len(part), "word_count": len(part.split()),
                "token_count": ntok(part),
                "hash": hashlib.md5(part.encode()).hexdigest(),
            })
            idx += 1
    for c in chunks:
        i = c["chunk_index"]
        c["total_chunks"] = len(chunks)
        c["prev_chunk_id"] = chunks[i - 1]["chunk_id"] if i else None
        c["next_chunk_id"] = chunks[i + 1]["chunk_id"] if i + 1 < len(chunks) else None
    return chunks


# ── pipeline ──────────────────────────────────────────────────────────
def sha256(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for b in iter(lambda: f.read(1 << 20), b""):
            h.update(b)
    return h.hexdigest()


def main():
    Path(OUTPUT_DIR).mkdir(parents=True, exist_ok=True)
    docs, all_chunks, errors, seen, seen_keys = [], [], [], {}, set()

    all_files = [p for p in sorted(Path(INPUT_DIR).rglob("*"))
                 if p.is_file() and p.suffix.lower() in EXTRACTORS]
    files = [p for p in all_files if not EXCLUDE.search(p.stem)]
    for p in all_files:
        if EXCLUDE.search(p.stem):
            errors.append({"file": str(p), "error": "excluded: template/master"})
    print(f"{len(files)} files ({len(all_files) - len(files)} excluded) in {INPUT_DIR}\n", file=sys.stderr)

    for path in files:
        try:
            doc_id = sha256(path)
            if doc_id in seen:
                errors.append({"file": str(path), "error": f"duplicate of {seen[doc_id]}"}); continue
            seen[doc_id] = path.name

            st = path.stat()
            fmt, pages = EXTRACTORS[path.suffix.lower()](path)
            pages = drop_boilerplate(pages)
            full = preprocess("\n\n".join(p["text"] for p in pages),
                              ocr=fmt.get("is_scanned", False))
            lang, conf = detect_lang(full)

            dm = {"doc_id": doc_id, "file_name": path.name, "file_path": str(path.resolve()),
                  "file_type": path.suffix.lower().lstrip("."),
                  "mime_type": mimetypes.guess_type(path)[0], "file_size_bytes": st.st_size,
                  "fs_modified": datetime.fromtimestamp(st.st_mtime, tz=timezone.utc).isoformat(),
                  "language": lang, "language_conf": conf,
                  "char_count": len(full), "word_count": len(full.split()), "token_count": ntok(full),
                  "is_empty": doc_is_empty(full),
                  "processed_at": datetime.now(timezone.utc).isoformat(), **fmt}

            if dm["is_empty"]:                       # empty doc must not emit chunks
                errors.append({"file": str(path), "error": f"no usable text ({dm['token_count']} tok)"})
                print(f"⊘ {path.name[:48]:48} skipped (empty)", file=sys.stderr)
                continue

            cs = chunk_document(dm, pages)
            cs = [c for c in cs if not (norm_key(c["text"]) in seen_keys
                                        or seen_keys.add(norm_key(c["text"])))]
            dm["chunk_count"] = len(cs)
            dm["languages_in_doc"] = dict(Counter(c["language"] for c in cs))
            docs.append(dm); all_chunks.extend(cs)
            print(f"✓ {path.name[:48]:48} {lang} | {len(cs):3d} chunks | ocr={fmt.get('ocr_pages',0)}", file=sys.stderr)
        except Exception as e:
            errors.append({"file": str(path), "error": repr(e)})
            print(f"✗ {path.name}: {e}", file=sys.stderr)

    json.dump(docs, open(f"{OUTPUT_DIR}/document_metadata.json", "w"),
              indent=2, ensure_ascii=False, default=str)
    with open(f"{OUTPUT_DIR}/chunks.jsonl", "w") as f:
        for c in all_chunks:
            f.write(json.dumps(c, ensure_ascii=False) + "\n")
    json.dump(errors, open(f"{OUTPUT_DIR}/errors.json", "w"), indent=2)

    tk = [c["token_count"] for c in all_chunks] or [0]
    print(f"\n{len(docs)} docs | {len(all_chunks)} chunks | {len(errors)} skipped/errors")
    print(f"tokens/chunk  mean={np.mean(tk):.0f} p50={np.percentile(tk,50):.0f} p95={np.percentile(tk,95):.0f}")
    print(f"langs: {dict(Counter(c['language'] for c in all_chunks))}")


if __name__ == "__main__":
    main()